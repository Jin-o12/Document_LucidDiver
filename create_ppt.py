# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # 16:9 와이드스크린 비율 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 테마 색상 정의 (Tailwind CSS 스타일 다크/틸 테마)
    BG_COLOR = RGBColor(15, 23, 42)       # Slate 900 (다크 네이비)
    TEXT_WHITE = RGBColor(255, 255, 255)  # 완전 화이트
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400 (연회색)
    POINT_COLOR = RGBColor(20, 184, 166)  # Teal 500 (포인트 틸/민트)
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800 (카드 배경)

    # 폰트명 정의
    FONT_NAME = "Malgun Gothic"

    def apply_solid_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, category, title):
        # 상단 카테고리 꼬리표
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = FONT_NAME
        p_cat.font.size = Pt(12)
        p_cat.font.bold = True
        p_cat.font.color.rgb = POINT_COLOR
        
        # 메인 타이틀
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p_title = title_tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_NAME
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # Slide 1: 타이틀 슬라이드
    # ----------------------------------------------------
    slide_1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_1, BG_COLOR)
    
    # 장식 데코 카드 (상단 틸 바)
    deco_shape = slide_1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(2.0), Inches(0.1)
    )
    deco_shape.fill.solid()
    deco_shape.fill.fore_color.rgb = POINT_COLOR
    deco_shape.line.fill.background()
    
    # 메인 타이틀
    title_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "수상한 손님들: 민원 폭발 대소동"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "Act Normal: Supermarket"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = POINT_COLOR
    
    # 하단 정보 및 설명
    desc_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8))
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    
    pd = desc_tf.paragraphs[0]
    pd.text = "AI NPC와 함께하는 3D 비대칭 소셜 디덕션 & 난투 추격전"
    pd.font.name = FONT_NAME
    pd.font.size = Pt(18)
    pd.font.color.rgb = TEXT_WHITE
    pd.space_after = Pt(20)
    
    pi = desc_tf.add_paragraph()
    pi.text = "유니티 3팀 기획 발표 자료  |  발표자: 장수영\n기술 스택: Unity 6 (URP) / PUN2 / OpenAI LLM"
    pi.font.name = FONT_NAME
    pi.font.size = Pt(14)
    pi.font.color.rgb = TEXT_MUTED
    
    # 대본 추가
    slide_1.notes_slide.notes_text_frame.text = (
        "안녕하세요. 3차 프로젝트 주제 선정을 위한 발표를 맡은 유니티 3팀 장수영입니다. "
        "저희가 준비한 기획은 AI NPC들이 가득한 마트를 무대로 한 비대칭 멀티플레이어 숨바꼭질 게임인 "
        "《수상한 손님들: 민원 폭발 대소동》입니다. 최신 생성형 AI 기술과 긴박한 피지컬 액션이 결합하여 "
        "포트폴리오적 가치와 대중적 흥행력(상품성)을 모두 확보한 독창적인 소셜 게임입니다. 지금부터 상세 내용을 소개해 드리겠습니다."
    )

    # ----------------------------------------------------
    # Slide 2: 기획 배경 및 의도
    # ----------------------------------------------------
    slide_2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_2, BG_COLOR)
    add_header(slide_2, "Concept & Intent", "왜 마트인가? 왜 2막 구조인가?")
    
    # 카드 1: 기획 의도
    card1 = slide_2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = POINT_COLOR
    card1_tf = card1.text_frame
    card1_tf.margin_left = Inches(0.3)
    card1_tf.margin_top = Inches(0.3)
    card1_tf.word_wrap = True
    
    p_c1_t = card1_tf.paragraphs[0]
    p_c1_t.text = "🎯 기획 의도 및 배경"
    p_c1_t.font.name = FONT_NAME
    p_c1_t.font.size = Pt(20)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = POINT_COLOR
    p_c1_t.space_after = Pt(20)
    
    p_c1_b = card1_tf.add_paragraph()
    p_c1_b.text = (
        "• 메타버스 요건의 재해석\n"
        "  - 거대하고 모호한 플랫폼을 배제하고, 플레이어와 AI가 긴밀히 교감하는 고밀도 3D 공간(마트)을 연동.\n\n"
        "• 2차 기술 자산의 마이그레이션\n"
        "  - 2차 프로젝트의 LLM 프롬프트 조립 엔진과 스트레스 연산 구조를 3D 멀티플레이어 환경으로 확장.\n\n"
        "• 제품이 아닌 '상품' 개발\n"
        "  - 기업과 시장이 좋아하는 캐주얼하고 유쾌한 상업적 공식 채택."
    )
    p_c1_b.font.name = FONT_NAME
    p_c1_b.font.size = Pt(14)
    p_c1_b.font.color.rgb = TEXT_WHITE
    
    # 카드 2: 핵심 메커니즘 개요
    card2 = slide_2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = TEXT_MUTED
    card2_tf = card2.text_frame
    card2_tf.margin_left = Inches(0.3)
    card2_tf.margin_top = Inches(0.3)
    card2_tf.word_wrap = True
    
    p_c2_t = card2_tf.paragraphs[0]
    p_c2_t.text = "🔄 공수교대의 기승전결 템포"
    p_c2_t.font.name = FONT_NAME
    p_c2_t.font.size = Pt(20)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = TEXT_WHITE
    p_c2_t.space_after = Pt(20)
    
    p_c2_b = card2_tf.add_paragraph()
    p_c2_b.text = (
        "• 1막: 영업 시간 (스텔스 & 수사)\n"
        "  - 도둑들은 숨고 경찰은 관찰하며 의심 대상을 검문함 (정적인 뇌지컬 영역).\n\n"
        "• 2막: 민원 폭발 (공수교대 추격)\n"
        "  - 민원 게이지가 폭발하면 경찰의 무기가 압수됨.\n"
        "  - 입을 막으려는 도둑들과 살아서 무전 지원을 부르려는 경찰의 개그 난투 추격전 (동적인 피지컬 영역).\n\n"
        "• 딜레이와 피로도 없는 깔끔한 액션 템포 보장."
    )
    p_c2_b.font.name = FONT_NAME
    p_c2_b.font.size = Pt(14)
    p_c2_b.font.color.rgb = TEXT_MUTED
    
    slide_2.notes_slide.notes_text_frame.text = (
        "이 슬라이드는 저희가 왜 이 게임을 기획했는지, 핵심 의도를 보여줍니다. "
        "저희는 교육 과정의 메타버스라는 조건을 모호한 친목 공간이 아니라, 다수의 AI NPC와 3D 멀티 환경에서 "
        "상태를 공유하는 알찬 공간으로 재정의했습니다. 또한 기존 2차 프로젝트의 기술을 버리지 않고 고스란히 확장하여 "
        "개발 효율을 극대화했으며, 뇌지컬로 시작해 피지컬 추격전으로 이어지는 1막과 2막의 롤체인지 템포를 설계하여 "
        "리플레이 가치가 매우 높은 상업적 콘텐츠를 타겟으로 삼았습니다."
    )

    # ----------------------------------------------------
    # Slide 3: 1막 - 영업 시간
    # ----------------------------------------------------
    slide_3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_3, BG_COLOR)
    add_header(slide_3, "Phase 1: Act Normal", "1막 — 영업 시간: 위장과 수사")
    
    # 3단 카드 구성
    col_width = Inches(3.6)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    steps = [
        ("👥 1. 도둑의 위장", "• 목표물 절도\n  - 마트 곳곳에 지정된 고가 물품 훔치기 미션.\n• 무빙 위장 (Space바)\n  - 플레이어 특유의 급회전을 보정하고 AI 손님처럼 일정한 속도와 회전각으로 걷는 의태 모드."),
        ("🚨 2. AI NPC 목격 핑", "• 실시간 감시망\n  - 도둑이 물건을 훔칠 때 3m 내에 AI 손님이 있으면 비명을 지르고 경찰에게 알림.\n• 간접 힌트 제공\n  - 범인을 직접 찍어주지 않고, 사건 현장의 대략적인 위치 핑만 경찰 미니맵에 제공."),
        ("🕵️ 3. 경찰의 관찰 & 검문", "• CCTV & 순찰 분석\n  - 목격 핑을 따라가 무빙이 수상한 용의자를 압축.\n• 텍스트 검문 시도\n  - 오사 타격 패널티를 피하기 위해 다가가 일관성 있게 텍스트 채팅으로 알리바이 취조.")
    ]
    
    for i, (col_title, col_body) in enumerate(steps):
        left_pos = start_left + i * (col_width + gap)
        card = slide_3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.0), col_width, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = POINT_COLOR if i == 0 else TEXT_MUTED
        card_tf = card.text_frame
        card_tf.margin_left = Inches(0.25)
        card_tf.margin_top = Inches(0.3)
        card_tf.word_wrap = True
        
        pt = card_tf.paragraphs[0]
        pt.text = col_title
        pt.font.name = FONT_NAME
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = POINT_COLOR if i == 0 else TEXT_WHITE
        pt.space_after = Pt(15)
        
        pb = card_tf.add_paragraph()
        pb.text = col_body
        pb.font.name = FONT_NAME
        pb.font.size = Pt(13)
        pb.font.color.rgb = TEXT_WHITE if i == 0 else TEXT_MUTED
        
    slide_3.notes_slide.notes_text_frame.text = (
        "1막은 본격적인 스텔스 추리 시간입니다. 도둑들은 목표 상품을 훔쳐야 하는데, "
        "일반 손님으로 완벽히 위장하기 위해 조작감을 AI와 통일하는 '의태 이동(Space 키)'을 쓸 수 있습니다. "
        "도둑이 훔칠 때 옆에 NPC가 있다면 비명을 지르고 소리 핑이 울려 경찰을 그 장소로 끌어들입니다. "
        "경찰은 주변 CCTV와 순찰을 통해 수상한 자를 좁혀내고, 무작정 공격하기보단 실수를 막기 위해 대화 검문을 진행합니다."
    )

    # ----------------------------------------------------
    # Slide 4: 2막 - 민원 폭발
    # ----------------------------------------------------
    slide_4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_4, BG_COLOR)
    add_header(slide_4, "Phase 2: Blackout & Chase", "2막 — 민원 폭발: 쫓기던 자들의 대역전극")
    
    # 좌측: 룰 전환 및 민원
    l_box = slide_4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    l_tf = l_box.text_frame
    l_tf.word_wrap = True
    
    pl_t = l_tf.paragraphs[0]
    pl_t.text = "📈 민원 게이지 & 무기 압수"
    pl_t.font.name = FONT_NAME
    pl_t.font.size = Pt(20)
    pl_t.font.bold = True
    pl_t.font.color.rgb = POINT_COLOR
    pl_t.space_after = Pt(15)
    
    pl_b = l_tf.add_paragraph()
    pl_b.text = (
        "• 실시간 민원 상승 (데드라인)\n"
        "  - 삼엄하게 감시하며 손님을 의심하는 태도 탓에 가만히 있어도 민원 게이지가 실시간으로 계속 누적됨.\n\n"
        "• 강력한 오사 패널티\n"
        "  - 무고한 AI 손님 오타격 시 공격권이 소모되고 민원 지수가 25% 급폭등.\n\n"
        "• 경찰 권한 해임 (100% 도달 시)\n"
        "  - 민원 폭발 즉시 경찰의 뿅망치(무기)가 압수되고 도망자 상태로 강제 전환."
    )
    pl_b.font.name = FONT_NAME
    pl_b.font.size = Pt(14)
    pl_b.font.color.rgb = TEXT_WHITE
    
    # 우측: 추격전 내용
    r_box = slide_4.shapes.add_textbox(Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5))
    r_tf = r_box.text_frame
    r_tf.word_wrap = True
    
    pr_t = r_tf.paragraphs[0]
    pr_t.text = "🏃 2막: 비상 탈출 추격전"
    pr_t.font.name = FONT_NAME
    pr_t.font.size = Pt(20)
    pr_t.font.bold = True
    pr_t.font.color.rgb = TEXT_WHITE
    pr_t.space_after = Pt(15)
    
    pr_b = r_tf.add_paragraph()
    pr_b.text = (
        "• 경찰의 목표: 도주 및 출구 생존 탈출\n"
        "  - 무기 없이 도둑들의 끈질긴 추격을 따돌리고 마트 출구(Exit)까지 무사히 뛰어서 탈출해야 함.\n\n"
        "• 도둑의 목표: 지원 무전 차단 (입막음)\n"
        "  - 경찰이 마트 밖으로 도망쳐 지원 병력(백업)을 부르기 전에 뿅망치/카트를 들고 쫓아가 격타시켜 눕혀야 함.\n\n"
        "• 피지컬 난장판 대격돌\n"
        "  - 2막 시작 즉시 복잡한 대화가 걷히고 유쾌하고 왁자지껄한 대환장 추격전 액션으로 전환."
    )
    pr_b.font.name = FONT_NAME
    pr_b.font.size = Pt(14)
    pr_b.font.color.rgb = TEXT_MUTED

    slide_4.notes_slide.notes_text_frame.text = (
        "2막은 전세 역전의 시간입니다. 경찰이 활보하는 것 자체로 마트 민원이 누적되며, 오사 시에는 민원이 급증합니다. "
        "이 게이지가 100%가 되면 즉시 경찰의 무기가 압수되고 강제 퇴거당하게 됩니다. "
        "이때부터 도둑들은 경찰이 밖으로 나가 지원군을 요청하기 전에 입을 막기 위해 뿅망치를 들고 경찰을 추격합니다. "
        "정적으로 머리를 굴리던 1막에서 탈출구를 향해 소리를 지르며 도망치고 쫓는 대환장 런 게임으로 분위기가 급반전됩니다."
    )

    # ----------------------------------------------------
    # Slide 5: 생성형 AI NPC 연동 & 딜레마
    # ----------------------------------------------------
    slide_5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_5, BG_COLOR)
    add_header(slide_5, "AI & Interaction Guard", "생성형 AI NPC: 대화의 명분과 기술 극복")
    
    card1 = slide_5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = POINT_COLOR
    card1_tf = card1.text_frame
    card1_tf.margin_left = Inches(0.3)
    card1_tf.margin_top = Inches(0.3)
    card1_tf.word_wrap = True
    
    p_c1_t = card1_tf.paragraphs[0]
    p_c1_t.text = "💬 검문 대화의 기능화 (Dilemma)"
    p_c1_t.font.name = FONT_NAME
    p_c1_t.font.size = Pt(20)
    p_c1_t.font.bold = True
    p_c1_t.font.color.rgb = POINT_COLOR
    p_c1_t.space_after = Pt(15)
    
    p_c1_b = card1_tf.add_paragraph()
    p_c1_b.text = (
        "• 대화의 시스템적 강제성 해소\n"
        "  - 섣부른 오사는 즉시 민원 폭등으로 이어지므로, 경찰은 활을 쏘기 전 최후의 확인 장치로 텍스트 검문 대화를 스스로 사용하게 됨.\n\n"
        "• 실시간 텍스트 일원화\n"
        "  - 리스크가 큰 실시간 음성(보이스)을 배제하고 오직 100% 텍스트 채팅으로만 검문을 작동하여 개발 안전성 및 완성도 확보."
    )
    p_c1_b.font.name = FONT_NAME
    p_c1_b.font.size = Pt(13)
    p_c1_b.font.color.rgb = TEXT_WHITE
    
    card2 = slide_5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = TEXT_MUTED
    card2_tf = card2.text_frame
    card2_tf.margin_left = Inches(0.3)
    card2_tf.margin_top = Inches(0.3)
    card2_tf.word_wrap = True
    
    p_c2_t = card2_tf.paragraphs[0]
    p_c2_t.text = "⏳ 지연 시간(Latency) 엄폐 설계"
    p_c2_t.font.name = FONT_NAME
    p_c2_t.font.size = Pt(20)
    p_c2_t.font.bold = True
    p_c2_t.font.color.rgb = TEXT_WHITE
    p_c2_t.space_after = Pt(15)
    
    p_c2_b = card2_tf.add_paragraph()
    p_c2_b.text = (
        "• 생각 중 말풍선 연출\n"
        "  - 질문 즉시 머리 위에 '음...?' 로딩 애니메이션을 띄워 API 연산 시간(2초)을 뇌정지 반응 시간으로 숨김.\n\n"
        "• 타이핑 출력 지연\n"
        "  - 텍스트가 서서히 타이핑되어 출력되게 속도를 제어하여, 도둑과 AI를 텍스트 출력 속도로 감별할 수 없게 차단.\n\n"
        "• 도둑의 '위장 답변 퀵 휠 UI'\n"
        "  - 도둑 유저는 마우스 클릭으로 간편하게 완벽한 AI 규격의 사과/핑계 텍스트를 송신해 은폐 가능."
    )
    p_c2_b.font.name = FONT_NAME
    p_c2_b.font.size = Pt(13)
    p_c2_b.font.color.rgb = TEXT_MUTED

    slide_5.notes_slide.notes_text_frame.text = (
        "저희의 기술적 엣지인 생성형 AI 파트입니다. 게임에서 대화를 강제하지 않고, 공격권 제한과 민원 리스크 때문에 "
        "유저 스스로 대화를 검증 수단으로 쓰도록 딜레마를 짰습니다. 또한 실시간 통신 딜레이 문제는 "
        "경찰이 질문하자마자 '생각하는 말풍선'을 띄워 자연스러운 연출로 승화했고, 타이핑 속도 지연 처리를 통해 렉을 엄폐했습니다. "
        "도둑 유저에게는 퀵 휠 UI를 제공해 손쉽게 완벽한 AI 양식의 대사로 둘러대게 보완했습니다."
    )

    # ----------------------------------------------------
    # Slide 6: 기술 구현 및 마일스톤
    # ----------------------------------------------------
    slide_6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_6, BG_COLOR)
    add_header(slide_6, "Technical Milestones", "기술 구현 및 개발 마일스톤")
    
    col_width = Inches(3.6)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    tech_steps = [
        ("🌐 PUN2 동기화", "• 룸 생성 및 입장\n• 역할 배정 FSM 관리\n• 플레이어 위치 & 애니메이션 동기화\n• 도둑 절도 상호작용 및 소리 핑 RPC 처리\n• 2막 전환 세션 데이터 백업 동기화"),
        ("🎨 URP 3D 공간 최적화", "• 단일 마트 씬 집중\n  - 저사양 빌드 타겟\n• 라이팅 베이킹 (Baked GI)\n  - 실시간 그림자 제거로 오버헤드 0화\n• SRP 배처 & 오클루전 컬링\n  - 드로우콜 최소화 및 고정 프레임 확보"),
        ("⚙️ AI 경량화 & 데이터화", "• 중앙 집중 LLM 매니저\n  - 1회성 API 호출 통합 관리\n• ScriptableObject 데이터\n  - 이름, 성격, 쇼핑 목적 등 페르소나 정보를 SO화하여 런타임에 군중에게 랜덤 믹싱 적용")
    ]
    
    for i, (col_title, col_body) in enumerate(tech_steps):
        left_pos = start_left + i * (col_width + gap)
        card = slide_6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.0), col_width, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = POINT_COLOR if i == 1 else TEXT_MUTED
        card_tf = card.text_frame
        card_tf.margin_left = Inches(0.25)
        card_tf.margin_top = Inches(0.3)
        card_tf.word_wrap = True
        
        pt = card_tf.paragraphs[0]
        pt.text = col_title
        pt.font.name = FONT_NAME
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = POINT_COLOR if i == 1 else TEXT_WHITE
        pt.space_after = Pt(15)
        
        pb = card_tf.add_paragraph()
        pb.text = col_body
        pb.font.name = FONT_NAME
        pb.font.size = Pt(12)
        pb.font.color.rgb = TEXT_WHITE if i == 1 else TEXT_MUTED
        
    slide_6.notes_slide.notes_text_frame.text = (
        "다음은 기술 사양입니다. 포톤 네트워크인 PUN2를 사용해 로비, 상태, 무브를 동기화하고, "
        "3D 마트 단일 씬을 배경으로 조명 베이킹과 SRP 배처 등 리깅 및 렌더링 최적화를 진행해 "
        "모바일과 저사양 PC에서도 원활히 돌게 할 것입니다. AI는 중앙 매니저를 두어 API 1회 호출 단위로 경량 통신시키고, "
        " ScriptableObject에 페르소나 프로필을 담아 런타임에 믹싱하는 정석적인 구조로 비용과 부하를 완전히 통제합니다."
    )

    # ----------------------------------------------------
    # Slide 7: 상업성 및 상품성 분석
    # ----------------------------------------------------
    slide_7 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_7, BG_COLOR)
    add_header(slide_7, "Commercial Value", "상업성 및 상품성: 왜 무조건 성공하는가?")
    
    l_box = slide_7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    l_tf = l_box.text_frame
    l_tf.word_wrap = True
    
    pl_t = l_tf.paragraphs[0]
    pl_t.text = "📈 확실한 인디 트렌드 타겟"
    pl_t.font.name = FONT_NAME
    pl_t.font.size = Pt(20)
    pl_t.font.bold = True
    pl_t.font.color.rgb = POINT_COLOR
    pl_t.space_after = Pt(15)
    
    pl_b = l_tf.add_paragraph()
    pl_b.text = (
        "• 입증된 흥행 장르 공식\n"
        "  - <Oh, Deer>, <Liar's Bar> 등 소셜 연기 및 의심 숨바꼭질 파티 게임은 현재 스팀 마켓에서 최고로 뜨거운 대세 장르.\n\n"
        "• 마케팅 0원: 방송 바이럴 최적화\n"
        "  - 1발 남은 절체절명의 검문 순간, 국어책 사슴/알바생 연기를 필사적으로 소환하는 예능감은 스트리머 합방 및 유튜브 밈 생성에 가장 적합한 1순위 구조."
    )
    pl_b.font.name = FONT_NAME
    pl_b.font.size = Pt(14)
    pl_b.font.color.rgb = TEXT_WHITE
    
    r_box = slide_7.shapes.add_textbox(Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5))
    r_tf = r_box.text_frame
    r_tf.word_wrap = True
    
    pr_t = r_tf.paragraphs[0]
    pr_t.text = "💰 확장성 및 수익 모델 (BM)"
    pr_t.font.name = FONT_NAME
    pr_t.font.size = Pt(20)
    pr_t.font.bold = True
    pr_t.font.color.rgb = TEXT_WHITE
    pr_t.space_after = Pt(15)
    
    pr_b = r_tf.add_paragraph()
    pr_b.text = (
        "• 맵 스킨 다각화\n"
        "  - 마트 외에 '백화점 명품관', '공항 면세점', '박물관 전시실' 등 다양한 탈출 맵 팩 확장 가능.\n\n"
        "• 코스튬 및 이모트 마켓\n"
        "  - AI와 구별되지 않으면서 개성을 뽐낼 수 있는 아바타 옷 입히기, 코믹 댄스 모션 스킨 판매 등의 강력한 소셜 부분유료화 연계.\n\n"
        "• 리플레이 가치 극대화\n"
        "  - 매 판 무작위로 바뀌는 NPC 페르소나와 돌발 이벤트로 리플레잉 가치가 무궁무진함."
    )
    pr_b.font.name = FONT_NAME
    pr_b.font.size = Pt(14)
    pr_b.font.color.rgb = TEXT_MUTED

    slide_7.notes_slide.notes_text_frame.text = (
        "저희 기획의 최대 무기는 상품성입니다. 최근 스팀 마켓의 인디 게임 트렌드는 "
        "유저가 직접 연기를 하며 다른 사람을 속이는 파티 게임들이 최고 동접을 기록하고 있습니다. "
        "저희는 K-마트의 유쾌한 땡땡이 감성을 얹어, 스트리머들이 짤을 양산하기에 최적화된 구조를 만들었습니다. "
        "나아가 백화점, 박물관 등의 맵 확장과 아바타 코스튬 의상 마켓 등을 통해 명확한 부분유료화 수익 모델까지 "
        "구상할 수 있는, 상업 가치가 매우 뚜렷한 포트폴리오입니다."
    )

    # ----------------------------------------------------
    # Slide 8: Q&A
    # ----------------------------------------------------
    slide_8 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_background(slide_8, BG_COLOR)
    
    deco_shape_8 = slide_8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(2.0), Inches(0.1)
    )
    deco_shape_8.fill.solid()
    deco_shape_8.fill.fore_color.rgb = POINT_COLOR
    deco_shape_8.line.fill.background()
    
    title_box_8 = slide_8.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.5))
    tf_8 = title_box_8.text_frame
    tf_8.word_wrap = True
    p1_8 = tf_8.paragraphs[0]
    p1_8.text = "Q & A"
    p1_8.font.name = FONT_NAME
    p1_8.font.size = Pt(54)
    p1_8.font.bold = True
    p1_8.font.color.rgb = TEXT_WHITE
    p1_8.space_after = Pt(10)
    
    p2_8 = tf_8.add_paragraph()
    p2_8.text = "경청해 주셔서 감사합니다. 질문이 있으시다면 말씀해 주십시오."
    p2_8.font.name = FONT_NAME
    p2_8.font.size = Pt(20)
    p2_8.font.color.rgb = POINT_COLOR
    
    slide_8.notes_slide.notes_text_frame.text = (
        "이상으로 저희 유니티 3팀의 3차 프로젝트 기획 제안서 발표를 마치겠습니다. "
        "질의응답 시간을 가질 예정이오니, 질문이 있으시면 편하게 말씀해 주시면 성실히 답변하겠습니다. 감사합니다."
    )

    # 파일 저장 경로 설정
    output_dir = r"C:\Users\jang9\OneDrive\Desktop\작업\3차 프로젝트\00.메인_기획서"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    output_path = os.path.join(output_dir, "3차_프로젝트_발표_PPT.pptx")
    prs.save(output_path)
    print(f"PPTX successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
