# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

# 콘솔 한글 깨짐 방지
sys.stdout.reconfigure(encoding='utf-8')

def extract_ppt_content():
    path = r"C:\Users\jang9\OneDrive\Desktop\작업\3차 프로젝트\00.메인_기획서\침몽도시_루시드다이버_기획제안서.pptx"
    prs = Presentation(path)
    
    for i, slide in enumerate(prs.slides):
        print(f"=== SLIDE {i+1} ===")
        # 텍스트 프레임 추출
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        print(f"  [Text] {text}")
        
        # 노트 추출
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print("  --- NOTES ---")
                print(f"  {notes}")
        print()

if __name__ == "__main__":
    extract_ppt_content()
